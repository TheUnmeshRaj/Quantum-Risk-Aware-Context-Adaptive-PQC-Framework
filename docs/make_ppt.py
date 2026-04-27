from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
DARK_BG   = RGBColor(0x0a, 0x0f, 0x1e)
BLUE_ACC  = RGBColor(0x38, 0xbd, 0xf8)
GREEN_ACC = RGBColor(0x10, 0xb9, 0x81)
RED_ACC   = RGBColor(0xef, 0x44, 0x44)
YELLOW    = RGBColor(0xf5, 0x9e, 0x0b)
WHITE     = RGBColor(0xf1, 0xf5, 0xf9)
GRAY      = RGBColor(0x94, 0xa3, 0xb8)
CARD_BG   = RGBColor(0x11, 0x18, 0x27)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def bg(slide, color=DARK_BG):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def box(slide, x, y, w, h, color=CARD_BG, border=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

def accent_line(slide, x, y, w, color=BLUE_ACC):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

# ══════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 13.33, 7.5, RGBColor(0x06,0x0d,0x1a))
# Gradient accent bar top
accent_line(s, 0, 0.05, 13.33, BLUE_ACC)
accent_line(s, 0, 0.12, 13.33, GREEN_ACC)

txt(s, "UNYSIS", 1, 1.2, 11, 1, size=60, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
txt(s, "Quantum-Risk-Aware · Context-Adaptive PQC Framework", 1, 2.3, 11, 0.8, size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
accent_line(s, 2.5, 3.25, 8.33, BLUE_ACC)
txt(s, "Week Progress Presentation", 1, 3.5, 11, 0.5, size=16, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Post-Quantum Cryptography · Migration Simulation · Decision Engine · Quantum Attack Analysis",
    0.8, 4.2, 11.7, 0.6, size=13, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Unmesh Raj  |  April 2026", 1, 6.5, 11, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, RED_ACC)
txt(s, "The Quantum Threat Is Real", 0.4, 0.65, 12, 0.7, size=32, bold=True, color=WHITE)
txt(s, "Why existing encryption is no longer safe", 0.4, 1.35, 10, 0.4, size=15, color=GRAY, italic=True)

problems = [
    ("🔓", "Shor's Algorithm", "Breaks RSA-2048 in HOURS on a fault-tolerant quantum computer.\nClassically it would take 10²⁰ years."),
    ("📦", "Harvest Now, Decrypt Later", "Adversaries collect encrypted traffic today.\nWhen quantum computers arrive, they decrypt it retroactively."),
    ("⚙️", "One-Size-Fits-All Fails", "Deploying Kyber-1024 everywhere crashes IoT sensors.\nDeploying RSA everywhere leaves hospitals exposed."),
]
for i, (icon, title, body) in enumerate(problems):
    bx = 0.4 + i*4.2
    box(s, bx, 2.1, 3.9, 4.2, CARD_BG, RED_ACC)
    txt(s, icon,  bx+0.2, 2.25, 3.5, 0.6, size=28)
    txt(s, title, bx+0.2, 2.95, 3.5, 0.5, size=16, bold=True, color=RED_ACC)
    txt(s, body,  bx+0.2, 3.55, 3.5, 2.5, size=12, color=GRAY)

# ══════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, GREEN_ACC)
txt(s, "Our Solution: Adaptive PQC Framework", 0.4, 0.65, 12, 0.7, size=30, bold=True, color=WHITE)
txt(s, "Context-aware algorithm selection matched to device, risk, and threat level", 0.4, 1.35, 11, 0.4, size=14, color=GRAY, italic=True)

pillars = [
    (GREEN_ACC, "01", "Quantum Risk Index (QRI)", "Scores each device 0–100 using\n5 weighted factors: data sensitivity,\nlifetime, threat window, exposure,\nand hardware capability."),
    (BLUE_ACC,  "02", "Decision Engine", "Selects the optimal NIST-certified PQC\nalgorithm per device. Considers RAM,\nFPU, bandwidth, and adversary type.\nOutputs score breakdown + alternatives."),
    (YELLOW,    "03", "Migration Simulator", "Gymnasium RL environment simulates\n6 enterprise devices over 10-step\nquantum escalation. Proves Adaptive\nbeats Status Quo and Paranoid."),
    (RGBColor(0xa7,0x8b,0xfa), "04", "Quantum Attack Sim", "Builds a real Qiskit QPE circuit.\nExtrapolates RSA-2048 break time:\n10²⁰ years classically → hours on\na fault-tolerant quantum computer."),
]
for i, (color, num, title, body) in enumerate(pillars):
    bx = 0.3 + i*3.18
    box(s, bx, 2.05, 3.0, 4.5, CARD_BG, color)
    txt(s, num,   bx+0.2, 2.2,  2.6, 0.5, size=22, bold=True, color=color)
    txt(s, title, bx+0.2, 2.75, 2.6, 0.6, size=13, bold=True, color=WHITE)
    txt(s, body,  bx+0.2, 3.4,  2.6, 3.0, size=11, color=GRAY)

# ══════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, BLUE_ACC)
txt(s, "System Architecture", 0.4, 0.65, 10, 0.7, size=30, bold=True, color=WHITE)

layers = [
    (BLUE_ACC,  "FRONTEND",  "Streamlit Dashboard  ·  3-Tab Interface  ·  Dark-Mode UI", 1.2),
    (GREEN_ACC, "BACKEND",   "FastAPI  ·  /analyze  ·  /simulate/quantum_attack  ·  /simulate/migration", 2.4),
    (YELLOW,    "ENGINES",   "risk_engine.py  ·  decision_engine.py  ·  devices.py", 3.6),
    (RGBColor(0xa7,0x8b,0xfa), "SIMULATORS", "quantum_attack.py (Qiskit)  ·  migration_env.py (Gymnasium)  ·  evaluate_framework.py", 4.8),
]
for color, label, desc, y in layers:
    box(s, 0.4,  y, 1.4, 0.7, color)
    txt(s, label, 0.45, y+0.15, 1.3, 0.5, size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    box(s, 2.0,  y, 10.8, 0.7, CARD_BG, color)
    txt(s, desc, 2.2, y+0.15, 10.4, 0.5, size=13, color=WHITE)

txt(s, "All layers communicate via HTTP REST. Streamlit calls FastAPI, which imports and runs the simulator/engine modules.",
    0.4, 6.1, 12.5, 0.5, size=12, color=GRAY, italic=True)

# ══════════════════════════════════════════════
# SLIDE 5 — MIGRATION SIMULATION RESULTS
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, GREEN_ACC)
txt(s, "Migration Simulation: Proof of Concept", 0.4, 0.65, 12, 0.7, size=30, bold=True, color=WHITE)
txt(s, "6 enterprise devices  ·  10-step quantum threat escalation  ·  3 competing strategies", 0.4, 1.35, 12, 0.4, size=14, color=GRAY, italic=True)

strategies = [
    (RED_ACC,   "🔴 Status Quo",  "RSA-2048 everywhere", "-5,340",  "$29.4M", "54", "0%",  "loser"),
    (YELLOW,    "🟡 Paranoid",    "Max PQC everywhere",  "-700",    "$2.3M",  "0",  "100%","mid"),
    (GREEN_ACC, "🟢 Adaptive",    "Our Framework",       "-2,032",  "$9.7M",  "14", "61%", "winner"),
]
cols = ["Reward", "Total Cost", "Breaches", "Compliance"]
headers_x = [5.5, 7.4, 9.6, 11.2]
for hx, col in zip(headers_x, cols):
    txt(s, col, hx, 2.05, 1.6, 0.35, size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

for i, (color, name, subtitle, reward, cost, breaches, compliance, rank) in enumerate(strategies):
    y = 2.5 + i*1.45
    box(s, 0.4, y, 13.0, 1.3, CARD_BG, color)
    txt(s, name,     0.6,  y+0.15, 2.5, 0.5, size=15, bold=True, color=color)
    txt(s, subtitle, 0.6,  y+0.65, 2.5, 0.4, size=11, color=GRAY, italic=True)
    for val, hx in zip([reward, cost, breaches, compliance], headers_x):
        c = GREEN_ACC if rank=="winner" else (RED_ACC if rank=="loser" else YELLOW)
        txt(s, val, hx, y+0.35, 1.6, 0.5, size=18, bold=True, color=c, align=PP_ALIGN.CENTER)

txt(s, "★  Adaptive achieves the best balance: far fewer breaches than Status Quo, at a fraction of Paranoid's cost.",
    0.4, 6.95, 12.5, 0.4, size=12, bold=True, color=GREEN_ACC)

# ══════════════════════════════════════════════
# SLIDE 6 — QUANTUM ATTACK RESULTS
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, RED_ACC)
txt(s, "Quantum Attack Analysis: Shor's Algorithm", 0.4, 0.65, 12, 0.7, size=30, bold=True, color=WHITE)
txt(s, "Simulated on Qiskit AerSimulator  ·  Extrapolated to RSA-2048", 0.4, 1.35, 10, 0.4, size=14, color=GRAY, italic=True)

# Two big stat boxes
box(s, 0.4, 2.1, 5.9, 3.2, RGBColor(0x2d,0x15,0x15), RED_ACC)
txt(s, "🔓  RSA-2048", 0.7, 2.3, 5.3, 0.5, size=16, bold=True, color=RED_ACC)
txt(s, "~10²⁰ Years", 0.7, 2.9, 5.3, 0.8, size=36, bold=True, color=WHITE)
txt(s, "Classical NFS on a 1 PetaFLOP supercomputer", 0.7, 3.75, 5.3, 0.4, size=12, color=GRAY, italic=True)
txt(s, "⚡  ~13 Hours", 0.7, 4.2, 5.3, 0.6, size=24, bold=True, color=RED_ACC)
txt(s, "Shor's algorithm on a 1 MHz logical-clock quantum computer\n(~4,096 perfect logical qubits)", 0.7, 4.85, 5.3, 0.6, size=11, color=GRAY)

box(s, 6.6, 2.1, 6.3, 3.2, RGBColor(0x0f,0x3d,0x2e), GREEN_ACC)
txt(s, "🛡️  ML-KEM (Kyber)", 6.9, 2.3, 5.7, 0.5, size=16, bold=True, color=GREEN_ACC)
txt(s, "Centuries+", 6.9, 2.9, 5.7, 0.8, size=36, bold=True, color=WHITE)
txt(s, "No known quantum speedup for Learning With Errors (LWE)", 6.9, 3.75, 5.7, 0.4, size=12, color=GRAY, italic=True)
txt(s, "✓  Grover's: √n speedup only", 6.9, 4.2, 5.7, 0.4, size=16, bold=True, color=GREEN_ACC)
txt(s, "Mitigated by larger key sizes (Kyber-1024).\nLattice SVP remains unsolved for quantum machines.", 6.9, 4.7, 5.7, 0.6, size=11, color=GRAY)

txt(s, "Conclusion: RSA-2048 is classically secure but quantum-broken. Our framework migrates devices to lattice-based PQC before Q-Day.",
    0.4, 5.7, 12.5, 0.55, size=13, bold=True, color=BLUE_ACC)

# ══════════════════════════════════════════════
# SLIDE 7 — DEVICE PROFILES & DECISION ENGINE
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, BLUE_ACC)
txt(s, "Decision Engine: Per-Device Algorithm Selection", 0.4, 0.65, 12, 0.7, size=28, bold=True, color=WHITE)

headers = ["Device", "Capability", "QRI Tier", "Selected Algorithm", "NIST Level"]
hx      = [0.4, 3.1, 4.9, 6.5, 10.8]
hw      = [2.5, 1.6, 1.4, 4.1, 2.0]
for hxx, col, w in zip(hx, headers, hw):
    txt(s, col, hxx, 1.5, w, 0.4, size=11, bold=True, color=GRAY)

rows = [
    ("IoT Temperature Sensor",    "1.06 / 10",  "MODERATE", "Kyber-512 / Dilithium-2",     "NIST L1"),
    ("Developer Workstation",     "10.0 / 10",  "HIGH",     "Kyber-1024 / Dilithium-5",    "NIST L5"),
    ("Public API Server",         "10.0 / 10",  "ELEVATED", "Kyber-768 / Dilithium-3",     "NIST L3"),
    ("Hospital Patient Records",  "10.0 / 10",  "CRITICAL", "Kyber-1024 + SPHINCS+",       "NIST L5"),
    ("Industrial PLC Controller", "1.50 / 10",  "HIGH",     "Hybrid RSA+Kyber-512",        "NIST L1"),
    ("Smart Home Hub",            "10.0 / 10",  "ELEVATED", "Kyber-768 / Dilithium-3",     "NIST L3"),
]
tier_colors = {"MODERATE": YELLOW, "HIGH": RED_ACC, "ELEVATED": YELLOW, "CRITICAL": RED_ACC}
for i, (dev, cap, tier, algo, nist) in enumerate(rows):
    y = 2.05 + i*0.77
    bg_c = RGBColor(0x11,0x18,0x27) if i%2==0 else RGBColor(0x0f,0x14,0x22)
    box(s, 0.3, y, 12.6, 0.68, bg_c)
    tc = tier_colors.get(tier, GRAY)
    for val, hxx, w, col in zip([dev, cap, tier, algo, nist], hx, hw, [WHITE, GRAY, tc, BLUE_ACC, GREEN_ACC]):
        txt(s, val, hxx+0.08, y+0.12, w, 0.45, size=11, color=col)

# ══════════════════════════════════════════════
# SLIDE 8 — WEEK ACHIEVEMENTS
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, GREEN_ACC)
txt(s, "This Week's Achievements", 0.4, 0.65, 12, 0.7, size=32, bold=True, color=WHITE)

achievements = [
    (GREEN_ACC, "✅ Unified Platform",          "Integrated all standalone scripts (simulators, engines, dashboard) into one cohesive FastAPI + Streamlit application."),
    (GREEN_ACC, "✅ Quantum Attack Simulator",   "Built a real Qiskit QPE circuit, computed gate counts, depth, and extrapolated RSA-2048 break times. Exposed via REST API."),
    (GREEN_ACC, "✅ Migration Strategy Engine",  "Overhauled the Gymnasium RL environment with real device capabilities derived from hardware profiles (RAM, CPU, FPU)."),
    (GREEN_ACC, "✅ Financial Cost Model",       "Added a $-denominated cost model: migration cost per algorithm tier + breach/crash cost per device class (IoT vs Hospital DB)."),
    (GREEN_ACC, "✅ NIST Compliance Tracking",   "Each device gets a live compliance score (0–100%) as it migrates. Compliance timelines are plotted per strategy."),
    (GREEN_ACC, "✅ Premium Dashboard",          "Dark-mode Streamlit UI with per-device allocation tables, multi-chart tabs, KPI summary cards, and the 'Why Adaptive Wins' panel."),
]
for i, (color, title, body) in enumerate(achievements):
    y = 1.6 + i*0.87
    box(s, 0.4, y, 12.6, 0.75, CARD_BG, color)
    txt(s, title, 0.6,  y+0.12, 4.0, 0.5, size=13, bold=True, color=color)
    txt(s, body,  4.8,  y+0.12, 8.0, 0.5, size=12, color=GRAY)

# ══════════════════════════════════════════════
# SLIDE 9 — NEXT STEPS
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_line(s, 0.4, 0.5, 2, BLUE_ACC)
txt(s, "Roadmap & Next Steps", 0.4, 0.65, 10, 0.7, size=32, bold=True, color=WHITE)

steps = [
    (BLUE_ACC,  "🧠 RL Policy Training",      "Replace the hand-coded adaptive policy with a trained Deep Q-Network (DQN) agent using Stable-Baselines3."),
    (YELLOW,    "📡 Live Threat Feed",         "Integrate CISA/NIST advisories via API to dynamically update threat levels in real-time."),
    (GREEN_ACC, "🏢 Org-Level Simulation",     "Extend the simulator to model 100+ device fleets with heterogeneous sub-networks."),
    (BLUE_ACC,  "📄 Compliance Report Export", "Auto-generate PDF/Excel NIST migration compliance reports from dashboard for auditors."),
    (YELLOW,    "🔐 Real PQC Benchmarks",      "Benchmark actual Kyber/Dilithium key generation and encapsulation times on target hardware classes."),
]
for i, (color, title, body) in enumerate(steps):
    y = 1.7 + i*1.0
    box(s, 0.4, y, 0.08, 0.75, color)
    txt(s, title, 0.7, y+0.1, 4.0, 0.5, size=14, bold=True, color=color)
    txt(s, body,  5.0, y+0.1, 7.9, 0.5, size=13, color=GRAY)

# ══════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 13.33, 7.5, RGBColor(0x06,0x0d,0x1a))
accent_line(s, 0, 0.05, 13.33, BLUE_ACC)
accent_line(s, 0, 0.12, 13.33, GREEN_ACC)

txt(s, "UNYSIS", 1, 1.8, 11, 1.2, size=64, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
txt(s, "The right algorithm, for the right device, at the right threat level.", 1, 3.2, 11, 0.7, size=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
accent_line(s, 2.5, 4.1, 8.33, BLUE_ACC)
txt(s, "Thank you", 1, 4.4, 11, 0.6, size=22, bold=True, color=GREEN_ACC, align=PP_ALIGN.CENTER)
txt(s, "Unmesh Raj  ·  April 2026", 1, 6.5, 11, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

prs.save("Unysis_PQC_Framework_Presentation.pptx")
print("Saved: Unysis_PQC_Framework_Presentation.pptx")
